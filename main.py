import os
import requests
from flask import Flask, request, jsonify, send_file
import google.generativeai as genai
from PIL import Image
import io
import base64
import tempfile
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json

app = Flask(__name__)

genai.configure(api_key=os.environ.get("GEMINI_API_KEY"))


# DGN Logo (base64-innebygd)
DGN_LOGO_B64 = "iVBORw0KGgoAAAANSUhEUgAAAR0AAAGQCAYAAACTRLR9AAABCGlDQ1BJQ0MgUHJvZmlsZQAAeJxjYGA8wQAELAYMDLl5JUVB7k4KEZFRCuwPGBiBEAwSk4sLGHADoKpv1yBqL+viUYcLcKakFicD6Q9ArFIEtBxopAiQLZIOYWuA2EkQtg2IXV5SUAJkB4DYRSFBzkB2CpCtkY7ETkJiJxcUgdT3ANk2uTmlyQh3M/Ck5oUGA2kOIJZhKGYIYnBncAL5H6IkfxEDg8VXBgbmCQixpJkMDNtbGRgkbiHEVBYwMPC3MDBsO48QQ4RJQWJRIliIBYiZ0tIYGD4tZ2DgjWRgEL7AwMAVDQsIHG5TALvNnSEfCNMZchhSgSKeDHkMyQx6QJYRgwGDIYMZAKbWPz9HbOBQAAAwXUlEQVR4nO3deXhU1d0H8N+5985MJpPMmkkmIYQkkISQsAbBsKSsylJlkaVifWm1CFGxKmq16IvF9kXbt4or1aJFa6utijtqRRFxARTZAgKSECIkIQvZM5lJ5p73DxnegCwhmXtmMnw/z3OexwXOOXeZ7z333nPvJQozZrO5jyzLLUTERRRFUXhGRsZiUcsH0N1Jwe4AAFxcEDoAIBRCBwCEQugAgFAIHQAQCqEDAEIhdABAKIQOAAiF0AEAoRA6ACAUQgcAhELoAIBQCB0AEAqhAwBCIXQAQCiEDgAIhdABAKEQOgAgFEIHAIRC6ACAUAgdABAKoQMAQiF0AEAohA4ACIXQAQChEDoAIBRCBwCEQugAgFAIHQAQCqEDAEIhdABAKIQOAAiF0AEAoRA6ACAUQgcAhELoAIBQCB0AEAqhAwBCIXQAQCiEDgAIhdABAKEQOgAgFEIHAIRC6ACAUAgdABAKoQMAQiF0AEAohA4ACIXQAQChEDoAIBRCBwCEQugAgFAIHQAQCqEDAEIhdABAKIQOAAiF0AEAoRA6ACAUQgcAhELoAIBQCB0AEAqhAwBCIXQAQCiEDgAIhdABAKEQOgAgFEIHAIRC6ACAUAgdABAKoQMAQiF0AEAohA4ACIXQAQChEDoAIBRCBwCEQugAgFAIHQAQCqEDAEIhdABAKIQOAAiF0AEAoRA6ACAUQgcAhELoAIBQiohGGGOkqqp+1apVppqamh5utztp7969VF5eTk6nk9LT08nhcNTX1tbuWbFiRZskSQ2ccxFdu2id2CZRN954o65v376ZxcXF1qKiIqqsrCSXy0Xp6ekUFxfXeOzYsd0rVqzwybJcr6pqsLsNF4hzrn/33XdNBQUFWS+//LK5f//+QxVFiSsuLia32002m42cTmfLli1bNk6dOtXjcDh2/Pa3vz3OGGvVqk9Mq4o551GPPPJI+ieffJJbVVU15tixY2nNzc0JTU1N0TqdLsLr9VJbWxspikI6nY6IiCRJqnS5XM2KonydkpKyYeTIkV8sWbJkF2PM19F2zWZzn6ampgKfz2fQatnaUxSFevfufcv+/fsfF9FeZ0mSRAcOHIhduXJl9r59+y6trKwc1tTUNOD48eNRROTknFNrayv5fD5SFIX0er3/r1YmJCS0MMa+TkxM3JCTk7N5+fLl+xlj9UFcHDiLEwcT28KFC/sWFhb+tKamZmplZWVCY2Oj0+12U2vrD1nCOSfOOTHGiDFGkiRRVFQUSZJ03OFwHE5JSfkqOTn5zQcffPALu91eG7KDAEVR6N577x2Yl5e3tF+/fntjY2PbjEYjZ4xxIrqgYjAYuMPhaO3fv//WmTNn5r/66qu9OtIHs9ncR5bllgttr7NFURSekZGxuPNrTVuc88jFixfPGTt27Mu9e/cuN5vNXFGUC94mjDH/NvFlZ2cfmDZt2v0rV65M8x8wOtk3PefcoEERMoIPJZIk0Ztvvpk1f/78v2RnZx+2Wq2qoihn3I5n2/bt/7ssy9xsNvPMzMxDv/zlL5986623+kpSYK7GBGSkwzk33HvvvbPWrVt39dGjRydWV1frfb4OD07Oy2AwUHx8/PHBgwf/bf78+Q9Pnz699Gx/FiOdHzz11FOx77333q8OHDgw/+jRo+lNTU0UyKOVXq8nu91en56e/mZ2dvbyVatWHbyQ+n/1q19dunv37jVutzugASHLMrlcrhffe++9+wNZLxHR7NmzLT6f7/rPP/+89NixY+5A1+8XGRnp+/nPf/7xM88809yRP//HP/6x94YNG+7at2/fNUeOHDH5RzOBotPpKCkpqWH06NFrVq5cucxqtdYEtIELIcsyLV26dG5ubu42m83GJUnSdFRhMBh4QkJC4fjx42/gnJ8xMC/2kc66deucY8eOvSc+Pr5Er9drvvyMMe50OmvGjh376IYNG1wd6SPn3DBkyJCvtOpTYmLiqgtfc+c3bty4ftHR0aeMGAJdiIhHRUXxX//611efrz+cc2nSpEnXJyQklMqyrPm2NhgMfMCAAXuWLVs2KXBr9QI8/PDDWcOHD19rtVrPODwLdGlfd2RkJB80aNBbjz/+ePrp/brYQuff//63zBgjzrnuzjvvvDM9Pf2UsBG1TXQ6Hc/MzDx45513Tpdl+az9ZYzRrFmzboyMjDzlxxuoPkmSxHv16vVE19bqmeXl5WWaTKbWQPX1XMswa9asZ87Vl2XLliVlZma+5l+PpPG2bl+cTqdn+vTpSznnYu5+c86lm266aXHPnj2rRaTr6cW/k0qSxJOSksoXLlw4o/1OfrGFDudc2rVrV9/Jkyevt1gswrdH+8IY43Fxcfyyyy77A+c86kz9feKJJ1wpKSkVWvahu4cOEfFBgwZ9cbZrU9dee216fHz8Dq3PLM5VoqKi+PDhw58tLy+P6/QKPRfOuXTiaMrmzp37ZExMDCcizU+nzlT8geNvu0ePHnz+/Pm3+i9yXWyhs3r16qHZ2dl7g7kD+ou/D1FRUXzKlClvFxQUnBI8jDEaN27cGlmWNTsqh0vouFyu0tdff916eh9+97vfTe3Tp0+Vf1lF9KV9aT8yNRgMfOLEiV8VFxfHd2W9ntGxY8eiOOfGUaNGPd5+OBcqxWw2q0OHDv01EZHBYOh9MYTOibsVM7Ozs6sCtSyBLDqdjqelpb39pz/9yeTv85IlS/IsFoum2yZcQsdisfAbbrhhVPv2V65ceUmfPn3qRbTf0SJJEh81atRXhw8fTujKuv0RzjmbNm3aEyaT6ZRRRigVk8mkTpo06caYmJh4WZbdotoNVug89dRTqYMGDaoO1HIEsviPhrIs84yMjLePHj0aqSgKZWdnfymi7XAIHZ1Ox6+88sq7/W1/+OGHjoyMjJ0i2r7QoigKHzNmzKcBnapwxRVXXO+/YByMIV1Hi91u982cOfMxRVEaRLUZjNApKSmxZ2VlbdJqmQJR/PuJTqfj48eP/3d+fv5So9EopN1wCB3GGB85cuR7siwT51zJzMz8WzCuoXa0yLLMR40a9URAgmflypXje/bsKWzk0NVitVq5Xq8XFo6iQ4dzLufl5T1lMBiCvq7PV/wjHpPJxOPi4gJ+p+psbYZD6BAR79Gjxz7OuTJt2rQrzGZz0Lfn+YrdbucLFiz4NWNdmP7HOY/OycnZE8qjmzMVkf0VHTo33HDDTIfDEfR13NHSPmhEbJdwCh2bzeZ5/vnnR2VlZRWIarOrpUePHjXLly/PPNd6POt9dkmSaOLEiTcWFBT0C9nnLs6iu/W3o/bs2dPr/fff/311dXWwu9Jh/md8/P8MHdfc3Kx//fXXny4sLMzq0uhBEEmSqKyszPrGG288xjk/6xMBZw2dP//5z8n79+9f6vF4tOkhXBBJkuiee+7JLysrO+dRBMIDY4y8Xi+tW7eun8fjoe4QOkREqqrSd999N2HevHlzz/Znzhg6jDF68cUXf11aWhrdXRY23L366quZu3btuinQz9VAaPKPEL1eL3HOqTu8VsTfx4aGBtq5c+cfzjZ/54yhc9dddyWXlJT8ojss6MVAlmV6+eWXlxw5cuSMs3wBQs2hQ4cSly9fnn+mQcuPQocxRlu2bPn18ePHre3PxyF4nnvuOcfnn39+ZVtbW7C7AtAhbrebPv/881+Ul5f/6ED5o9D57LPPeh05cuQXgXw1BXQeY4zeeuutX1VUVDiD3ReAjmKM0eHDhxNvuOGGGaf/vx+FzqOPPjrp6NGjViE9g/NSVVW/b9++a3EtB7oL/ylVS0sLKykpuen0F72dEjqcc7mgoOBnLS0t4noI53TdddcNLikp6RvsfgB0VPvLMiUlJUMeeOCBEe3//ymh88wzzwysra3Nw3Wc0MAYo8rKykVNTU1nf0ENQAirq6vTffnll6fcPj8ZOowx2rhx47TKykp8liZEqKpqKyoqmoC7iNBdtbW10aFDhyZzzk++beBkwKiqyvbs2TMG1w5Cx/Lly4dUVFT0CHY/ALqipKQk6bbbbkv2//vJ0HnllVeSKyoqBgWjU3BmW7duHV1TU4PZmdCtNTc3G8rKyn7q//eTofPxxx8Pb2xsNAenW3A6nU5HhYWFozB1Abo7r9dLhYWFo/3vVD4ZOt99992Ypqam4PUMTvHaa6/FHD9+vH+w+wEQCLW1tYNra2vNRCdCh3Ouq6mpyQ32BUtJkrrNg21ae+ONNwa63e7YYPbB//XH0/8d2yg0nG37dOTPilZfXx//9NNPpxOd+Jb5K6+84jhy5Ehg33HaCaqqksFgIL1e3+Jyuaqqq6u3x8bGqo2NjWSz2fqXlZXF1tbWRl0MjwMUFxcPDfYT/pxzUhSFLBYL2e32iuPHj39jt9s9Xq+XIiMj0+rq6hKqqqqsXq83qP28WHHOyWQykclkOm4wGIpbWlq+j4uLo5aWFqqrq9NZLJYhFRUVjubmZl2wfzN1dXWsqKhoEBFtVYiI3n333d5erzfmxBcfgtIpvV5PSUlJh5OTk1f17t377b/85S9HZFmur6mpIc45ff/991H33ntv3JYtW2YUFRXdXVRU5CD64WFIVVXD6hkxSZKourp6aDDuJPr3AcYY9ezZsyopKenZ8ePH/+f666/fkZqaevz48eP+yV8Rjz32WMz27duv3bZtW/7+/ft7tra2+r+lLbzfFxun08n79OnzYVJS0tM/+clPPsvPz69hjLVWVVUR0Q/70NatW60PPfRQyoEDB2YWFxcvLCkpcba1tVEwfudtbW20e/fuXif/w89+9rMlIl/xeXqJjo5uGT169P9+/fXXMefrPGOMVqxYkTxmzJhPjUYjlyQpqO9t1uLNgZxzQ8+ePXcEa5lMJhOfMWPGxg8++CClI0Pyr7/+OmnixIn/MJlMPNjv8Q2nNweeqUiSxHNycsp/85vfXBEREdGhfjPG6P333+972WWXvRMVFRW030tubu4HOp3uhw6lp6f/KVg/XofDoebl5c270PNNznnsJZdc8s6ZPhIvsmgROqtWreqRmJjYGIzliYiI4Dk5OXef681vZ8I5Z2PGjAn6u5vDOXQYYzwpKenLF154IaUz/eecG2bNmvWviIiIk/WJ7H9SUtI3nHOF9Ho9XXLJJZ8FYwUajUa+aNGix05/IKyjNmzYEJWSkhLU98dqETp33HFHblRUlEfkckiSxHU6HZ8yZcpLitK5F/pzzvVZWVnvnf5tbpElnEMnNTW1ZtWqVf26sgyc84gRI0ZsDcanpPr06dOwY8eOdOKcS1arVbOP2Z+tMMb4iBEjNnX2e8icc8Y5t82fP3+8zWZrE91/f9EidK666qqfR0VFCV+Wvn377tu+fbu1s/3mnNsfeuihjISEhHqi4H2BMhxDx2Aw8P/6r/+6LRDL8eCDD17qcrm8opfBZrM1z507N1167rnnko1GY3ogFqYj/KdRFouledq0ab9ljHXqquOJHbp2zZo1H2dlZX0UTrdxS0pKejQ3Nwtt02Qy8dTU1CWDBw+u7WwdjLHj99xzz/6cnJy/d3b0Cqfy79fJyclFzz///HOBqPPee+/dnJqa+qb/U9yiNDc3y4yxRGnjxo0x9fX1F3T+3lWMMerRo8e7v/nNbzZ1sR7OGOO5ubkPWq3WsHloLCIiIkdVVaHzKtLS0grffffdT7paj6qqNHv27EdiY2PrwumOYjDpdDqKi4t7lDFWF4j62traaODAgX89MZoWRqfT6c1m82CpV69eIyRJEhY6nHMym800YcKE1wJV5x//+MeNTqdzW6DqCybOOfv222/9kzaFtKnX6yk9Pf15xlhApqTPnz//oMvleiecRp/BwjmnmJgYPmPGjE8DWe8dd9yxxeFwlASyzvPxer20ceNGr7Rx40ZV9OSumJiYysWLF28IVH2MMTU7O/tNvV4fqCqDyRQfHz9UZIM2m803duzYtwNVH+echg0bttZkMp3/D8N52e32vbfeeuu+QNbZp0+fuvj4+K8CWee5MMaotbWVrFZrqlRWVuYQHTqRkZEfp6enVwSyzn79+r0dGRnZ7T/S9eyzz7Kqqiph6ckYI5PJtG3RokV7Alnv5MmTP42Oji4LZJ0XI0mSyGw2f8kYC+jrPDnnFBcX97nBIO7KCuecEhISLpEyMzPHijz3NhgMZLfbPwv0rNUHHnhgf0xMzMGAVhoEH3/8sSLydRaMMbLZbF8xxgI6T3769OlVLpdrm78N6BxZlik+Pv4bLepOSUnZrSiK0OnjpaWlPqmkpEToxSRJknwJCQk7NKi3zW63f9ndd/Dx48cPsFgswh70lGWZbDbblkDXq6oqxcTEfKrT6cLqERXRoqOjKTU19ZAWdfft2/d7m80mJHT8+0BFRQVJhw8fVkX+UO12Oxs2bFhloOvlnJMkSZ90dmJbqNi+fbvO6/UKu5cZHR1NWVlZAT3V9ZMk6TNFUYTuX+Gmra2tVpblvVrUvWDBghKPxyPs7ODEO79VSafTCX19Qmtra8n8+fPLtah76tSpx6Kjo7v104Y7duzgIp8ub2trq09JSSnSou4FCxZUxcTEtGKk03lWq7Vt9OjRDVrUrdPpWmJiYoRNCDvxVHys1LNnz34idwqj0VjtdDoDMt/gdG1tbbskSdKkblF2797tE/l0udFodNfU1BzTou6rrrrqkKqqhVrUfbFwOp00depUTYaKOp2OkpKShAxD201y7Cc1NGgSomeVkJAga3UKdP/997stFstxTSoX5LLLLssS+WoIh8Mh3X///ZqczhkMhjaXy6VJoF0sjh49upOINDmQejwe2r9//xciTn/9Axu3202S2+3WvMH2oqKiWrSafs0Ya6irqyvQpHJBXC5Xmsjp6Uaj0UNEmqSc1+ul0tLSr0VPtw8niqI0MsY0eVG2qqqk1+vrRG4fr9dLkui303322WefafUFUcYYZWdnd+sP09XU1PhEnu7u2rVrGxFpNtwdOnSojNDpvPj4eKblSKRXr16SyAv9ra2tJImcGMgYI0VRNL2wGB0drcmdGFEaGhqE3mLW6XRtJx6e1URcXFy33h7BZrVaNa3f6XQKnUfl8/lIEv3uVIPBoNkScs5p69atX8hy9x3siD7d1XJ7EBG9//77m3HLvPMiIyM1DQWHw6FZ3WeiqipJor+rpPU8mn79+kndea6O6Pcia33qM2TIEKmjr9WEH4uKitI0sSMjI7nIg4KqqiSF20u0U1JSqDuPdMLtywpZWVkk8vmecCPLsqY7xCuvvLJJ5NmOqqokiZ64pfWRNT4+vluHDpG4V1oQaf9cVGJiouaj23D2xhtvbNByYFBUVNTg//qHCJxzCrvbCna7XcXdktCRmJjIEDqdV19fr+ntZavVKkmSJPRAJ/R2mQhr1qzZ2tjY2G3PUcLtK6d6vb6VNJoHFM78X+TU6XSa7gwRERGan32cLuyGBLt37272+Xx42CdE3H333btLS0sxKzlEBWMUGnYjHbPZzHB6FTp27tzZyjnHSOcC+U93tN6XFUURPrLGrxM05XQ6NZ1RC12j0+mEho4kSQgdABArLEMHR9bQgu3ROSJvY4vCGAvP0AGA0IXQAQChEDoAIBRCBwCEQugAgFAIHQAQCqEDAEIhdABAKIQOAAiF0AEAoRA6ACAUQgcAhELoAIBQCB0AEAqhAwBCIXQAQCiEDgAIhdABAKEQOgAgFEIHAIRC6ACAUAgdABAKoQMAQiF0AEAohA4ACIXQAQChEDoAIBRCBwCEQugAgFAIHQAQCqEDAEIhdABAKIQOAAiF0AEAoRA6ACAUQgcAhELoAIBQCB0AEAqhAwBCIXQAQCiEDgAIhdABAKEQOgAgFEIHAIRC6ACAUAgdABAKoQMAQiF0AEAohA4ACIXQAQChEDoAIBRCBwCEQugAgFAIHQAQCqEDAEIhdABAKIQOAAiF0AEAoRA6ACAUQgcAhELoAIBQCB0AEAqhAwBCIXQAQCiEDgAIhdABAKEQOgAgFEIHAIRC6ACAUAgdABAKoQMAQiF0AEAohA4ACIXQAQChEDoAIBRCBwCEQugAgFAIHQAQCqEDAEIhdABAKIQOAAiF0AEAoRA6ACAUQgcAhELoAIBQCB0AEAqhAwBCIXQAQCiEDgAIhdABAKEQOgAgFEIHAIRC6ACAUAgdABAKoQMAQiF0AEAohA4ACIXQAQChEDoAIBRCBwCEQugAgFAIHQAQCqEDAEIhdABAKIQOAAiF0AEAoRA6ACAUQgcAhELoAIBQCB0AEAqhAwBChWXocM6D3QVoB9sD2gvL0AGA0BV2oePxeHBkDSGVlZXB7gKEGCXYHQg0j8fjC3Yf4BSqqqo4CsBJYTfSueWWW1IjIiLkYPcDfvDcc8/1sFgs5mD3A0JH2IVOampqWkRERNiN4LqrmJiYJIQOtBd2oVNZWelTVTXY3YATysrKVK/XG+xuQAgRHjptbW2a1l9SUkI+X/e9rCNJYjeJ1hfdDx48SK2trZq2Ad2L8NDxeDya1n/o0CHNgy2caD0K2bNnj+bbHLoX4aFTX1+v6TBk3759CJ0L4Ha7Nd0eO3fupJaWFi2bgG5GYowJa4xzToMHD+6r0+k0qZ8xRiNGjBiN06uO69WrVyrnPEKr+i+//PKRWtUN3ZPQ0CEicjgcaYqizc0lzjkdPXrUqknlgojeHmazOZmIjFrVX1xcbMFkTWhPEX1kLS0t9Wl1+sM5l10ul2Y/IBG0CuSzqamp8W3btk2T232SJFFRUVEU7iZCe5Lonby0tFTyer2aJF1xcXF0VFTUQC3qFsVoFJuZlZWV7NFHH9VkJ1AUhaxW63CMdKA9SZbFTt41m83p69evT9ai7pdeeim2uro6QvQpSiDp9Xqhp1hGo9ExaNCgbC3q3r9/v7WsrMyuRd3QfUlaXdQ9m7q6Ov3777/v0KLuXbt29VVVNbo7H1kjIyOFttfc3Cxt3rw5Tou6n3zyySRVVZO0qBu6L0mv1wtt0Ov1yg0NDTla1F1VVZXb2NioRdXCWK1WWeRIx+PxkKqqI7Sou7i4OKehoYF155EnBJ7w0HG73VRUVDQ00Dsi51yuqKjI6+4XLWtqag6LXAafz0clJSUDOecB3SCSJNGxY8fGejwe4XfkILRJoi9cnritPVxV1YDODfnnP/+ZWl1dPTSQdQbDCy+8sF306eGxY8f6f/TRR7GBrNPn81mqqqomcM6pux8IILAkj8dTLbrR48eP93388ccHBLLO1157La+qqkrssE0DmZmZsujrbDU1NY633357XCDrvOWWW4Z8//33zkDWCeFBKi4u3i2yQcYYHT9+XFm7du30QNXJOWeHDh36VTg8zdyvXz8SecorSRI1NjbSpk2brgrUnC3GGO3atev65uZmvGIEfkSKiYkROjuQc06tra20d+/eWf/4xz9sgajz9ttvzysuLs7pznet/NLS0pjI0OGcE+ecCgsLxy1ZsiQtEHWuXr069bvvvpuK0yo4EykxMTEoDdfW1qa9+eab93a1Hs65bsOGDb+tra3VhcMFy+zsbK9OpxP2a/UHdWNjo23nzp3/G4h5W6tXr761oqLC2uWKICxJSUnBmUbh9Xppw4YNC1asWJHX2ToYY7Ro0aIFxcXFl/mP2N3dhx9+uKuurq5CdLs+n4+2bNky5cYbb7ymK+F9++23j9+7d+8v8aQ/nI1UUFDwsejnr4h+CIyqqqro1atXPz1jxoxO3TnJz8/Peuuttx6ora0NcO+CZ8aMGT6HwyE8PSVJovr6euX1119/dPz48ZmdqeOll17q+fHHH/+zrq4uKtD9g/AhxcXF1Ym+W3KycUmiQ4cO9d2+fftH99xzz086+vdkWaZly5Zdtnbt2nfKysrs4XBa5Tdt2jRutVqFv/WKc06MMTp69Khj9+7d7y9evHjOhazXW2+9NfXOO+98ZefOnQG99Q7hRxoyZEiDoijCr/hxzsnn85GqqnT48OHsv/71r2+OGjVqxS233JLEOT/j0EuSJFq3bp1zzpw5j/ztb397p7y8PNlfVxhpKi8v/1p0o+3n01RUVCT9/e9/f/GSSy555sYbb8w628RBxhht3rzZPGnSpBteffXVL44cOYKHO+G8lM2bN3/OOfeQhu9U6YiqqirLli1b7i4oKPjV9u3bj15zzTVfFBcXf+Lz+Y5XV1dTbm7uyD179qQtWLBgbFVVVUK4vgKTMcbHjBnDP/nkk2D2gerq6nTffPPNgv3791/96aeffvvTn/70CyLaWV5e/r3H46EBAwYM/v7779Nnz549obq6Orm5uTlo/YXuRVmyZAndfvvtrKSkJGid8B8dW1tbqba2NmbTpk0xkiQNJKJ8g8FAPp+PioqKSFXVk3+WMRZuI5yTGhsbv5IkaXYwLo63b8/n81FdXV1UXV3dJd9+++0lREQ6nY4457R3795TRkfhvD0gsKSrrrrqMOe8KNgdOZ2qqqSqKrndbvJ6veTz+U7ZqcN5B+/Tp0+Z0WgM2jL6w+70APL5fNTS0kIej+fkqXH7vwPQEZKiKO6oqKjCYHcE/l9qaupBSZLw3RYIS5LP5yPG2E5JkvA0cIjIyckpio6Obgh2PwC0IBERDRgw4GCwbpvDj82cObMyMjJyX7D7AaAFiYioX79+BWazWcV5eWiQZZlbrdadGHlCOJKIiO6777790dHRR4LdGfiBqqrUt2/fPQaDIdhdAQg4/yS8ppiYmG04soaO4cOHf2mz2TD0hLAjEf0wIS0tLe1T0a8uhbO7+eabv7VarfuD3Q+AQDv5uMG4ceO+dDgceDQ4RDDG3BEREe9g9Anh5mToXHfddTsdDofQtwjCuY0YMWK91WoNdjcAAupk6DDGWhISEl7GKVboWLFixTcJCQnC360DoKVTnuaeOXPmeqvV2kJE3XaiYGRkpNB3DGvJarVWOp3Of4n+CiuAlk4Jnfz8/G9iY2PfCcZLvQJl0qRJe10uV1jM5lVVlWbOnPkvi8USno/Uw0XplHRRVZXy8vJeNJlMwepPp/hHZU6ns2nevHk3tLW1fRvkLgXMbbfd9rnL5frQv4zddQQK4PejIc2TTz65zul0bu5Os5MZY6TT6SgjI2P13LlzP7fb7WHzAKvP56MZM2a8aLfbg90VgID4Uegwxlpzc3P/GB0dHYz+dIqqqpSSklJ93333Perz+cjhcOwKp+sgv//9799wOp3vY5QD4eCMF29efPHF9Tk5OfuIQns477/2FBUVRbm5ufdcfvnlh4iITCbTloiICGKMad5/zjlp/RZDxphnwIABvzebzW3daQRK9MO2GThwIK5JwUlnDB3GWMPkyZOXOJ1Obyjv5JxzkmWZhgwZ8s6aNWue9f/3mTNnFplMploRfeecU1lZmU/rdv71r399kZub+053G8FlZGRstNvt7we7HxA6znqb6q677vpw4MCBT+r1eiEjhs7gnFNaWlrJsmXLbmWMnXyN3fXXX1/qcrnK/X9G6z707t27h9brhzHG8/Pzb09OTq7UtKEA8I9AY2NjW+6+++77Kysra4LcJQghZw0dxljr0qVL/yclJeWbE/8ecsGTmJjYNnHixKvHjx9/yoVjSZJa7Xb7VhH95ZxTUlLSMBHTDKZNm3Zo4sSJd4f6g6Ccc9Lr9TRy5MhHZs2atUWSpPCYOAUBcc5fytixY6tuvvnmhcnJybWh9l1qi8XiGzduXP5TTz31xen/j3NOmZmZH0ZGRgrpS0FBQavPp/kZFhERPfPMM88NGDDgf4zGoH6845wkSaJBgwa9s3bt2nsZY24WakcrCKrzHp4XL178dXJy8ly73e71f5At2PtQVFSULz09Pf+FF15YfbYf+7Bhw9ZHR0cf17ovjDFSFCWNcy7kdp+qqvTJJ5+s6N+//1uKooTE9iD6YT1IkkSyLFNiYuKnkydPvsZ/yhvK1wUhREmSRPfdd9+9PXr04EQU1GI2m9uGDx++oCN9Hjx48OuMMc37FBMT07hixYrUC1mnXfXBBx+YMjMz39TpdEHfJv6iKAofOHDg7n//+9/92vd1wIAB/xTVB8YY79Wr1xNdWbdnk5eXl2kymVpFLYskSdxgMJx3X++KtLS0PL1eL2wfSUxMPPOXNE+nqio9+OCDv7/++uvvj42N9RH9/5FNa+2P5ImJic1z5szJ37p161/P9/dUVaWRI0euNpvNmo8GPB6PqaSkJEezBs7g8ssvb3rqqafm9+3b913/iCdY36QnIjIYDJSdnf3xihUr5s6ZM2ev8I5AeJJlmYYNGzY1Ojq6lk4kMQlIR51Ox/v27btv0aJFl17ID4tzbhswYMA6xhjXcsQjSRKfNGnSc6JPczjnxoKCAvugQYPeFHm0Or3Y7XZ+5ZVXruGcm8/UT4x0Or9fheNIp1MdnTRp0mXJyclVge6QPxzaB4TT6fSNHDnyybVr1zo609enn36694ABA6rOFDqBDKKsrKzvOeeWzvSxqzjnpiuuuOJNi8US8OU6vbSv22g08sGDBx9cvHjxLEVRzto/hE7niqjQMRgMQpaHqAuhQ0S0fPnyjKFDh35pMpk02aHtdjvPycn5fOHChWPOtUN3xN///vfxOTk56+12O1cURZOVaTKZ+NVXXz29Sx3tAs65YdGiRb/t1avXcVmWNf1RG41Gnp6efmD27Nm3cc6t5+tbuITOqFGj+hmNRmE/UCLiiqIs1GJZ/Hr37j1G5HXBHj168E7/mv/7v/97P+d8/M9+9rPr9uzZs+zgwYMxLS0tna2OiH44fbPZbG0xMTEf5Obmrn722Wc/YIy5u1QpEV177bUfcc4/veuuuyYWFhZO2bJlS0zv3r2H1tXVUWtrYD6kyRhjkiRNkSTpjWBML2CMeTjnf37++effWbVq1d3ffffdzJqaGkMg7hxJkkRGo5Gio6Or4uLiPk1JSXlh1apV6+Pj45tC4c6ZKKmpqR63233A7XYLmRYuyzIpilK7fft2zdqYO3eu6YMPPih0u7v8M+uQ+Ph4Csges3Xr1p6/+93vJhcWFl537NixIU1NTTqv13vev8cYo8jISIqMjPQ6HI49cXFxb48fP/6tBx54YFugwuBMZFmmtra2gH/fpbi4mKWkpHQteQNAkiRaunTpJZs3b7714MGDY6uqquKbm5tJVdVTbl8zxn50O5sxRkajkRhj3Ol0eqKjo7+1Wq1fZmRkfHHppZduWLhwYemFzkkaMGDAP3ft2nV1QBbuPBhjlJSU9OThw4dvDnTdnHNGRKInOrYxxjSbBBaMZQroYYpzLt1+++39v/rqq5GNjY1jamtr+xBRYk1NDbndblIUheLi4qitra3IbreXWyyWA7169dqTlpa2+b777itkjOHF8AEkSRJ98sknzjVr1uQVFxePLC4uzoyIiBhcUVEhNTY2kqqqZDKZKD4+nmpqanampqY2GQyGQovF8q3ZbC644oorKq666qrDkiT5ujJiCpfQgRB34iiq27Rpk23evHm2pKQk26hRo2ybNm2ycc6li2lYHioURSHOufXuu++29e/f35aUlGTLz8+3cc5tOp1Os6kF4XJNBwC6CYQOtNd9X4YMAN0SQgcAhELoAIBQCB0AEAqhAwBCIXQAQCiEDgAIhdABAKEQOgAgFEIHAIRC6ACAUAgdABAKoQMAQiF0AEAohA4ACIXQAQChEDoAIBRCBwCEQugAgFBd+4odBFx+fn7md999l9KRT/gEgsvlol/84hdfTZkypVJIg3DRQ+iEmMLCwkVbtmy5pa1NzNd4EhISaPbs2VOI6D0hDcJFD6ETYtxud1tTUxOJ+kpoU1MTeb1e8Z8khYsWrumEIJHfBAvH748F47PO0HEIHQAQCqEDYUfU9TDoHIQOhJ3W1tZgdwHOAaEDYcfn8wW7C3AOCB0IOx6PJ9hdgHNA6EDYwd2r0IbQgbCDazqhDaEDYYVzjpFOiEPoAIBQCB0AEAqhAwBCIXRAc6KvsUgSdutQhq0DmuOcC2uLMUYGg0FYe3DhEDqgOdHPQiF0QhtCBzTX2NgotL2IiAih7cGFQeiA5qqqqoRe1ImKihLZHFwghM5FTlVVampq0qx+znlE7969+2nWwGkYY2S320U1B52A0LnIcc6pvr5eyyYUSZJcWjZwOofDIbI5uEAInYtca2srlZaWalb/woUL1aNHjwp714QkSdSjRw9ZVHtw4RA6IUb0HBOv10s7duzQ7JpLRkZGGudc2NBDVVXav3//V6LagwuH0AkxRqNRaHsej4c8Ho9Nq/r/85//ONva2oQtlKIodOjQoSOi2oMLh9AJMaJDh4goMTExV6u6vV5vbktLi1bV/4jRaKQxY8bg00ohDKETYhwOhyzyFEtVVSoqKkrinAf8WzQ6nY5qa2tHi3x9qKIovqFDh7qFNQgXDKETYqqqqvaJfFZJVVWqr6/vTxp8ePEf//iHvbS0dECg6z2X5ubmmrfeemuXyDYBujWLxTJGURRORMKK3W73/OEPfwh4OMyaNeuayMhIzhgTtix9+vSpOnz4sGbXqKDrMNIJMRMmTJBFX9dpaGjQf/TRR1cEsk5JkujgwYPXNzc3C33gU5blo0lJSXgzewhD6ISYyy67TI2MjBT3K6Uf5uocPHjwF5999ll0oOpcsGDBxAMHDowIVH0dZTQaixhjzaLbhY5D6ISYQ4cO7WpqaqoQ3W55eXmfJ5544v5AXFDmnFu3bdv2qNvtFvq4N2OMZFneKbJNgG6Pcx6dkZFRRgKv6fhLz5491T//+c83Mtb53OGcSxMnTnzGYDAI7TtjjEdGRvJZs2b9stOdB7gYcc71mZmZn5HgwPFf7E1NTfUsWbJkvk6n60zfI6ZPn/6E2Ww+pU5R/bdYLK0333zzwAvuOMDFLi8v7yXGmNAfbftit9t9o0aN+tuyZcv6deSFWHq9nh5++OGRo0eP3mQ0GoPSZyLiiYmJZV9//XVMR9czBEfAJ4RB102ZMuW29evXP9za2ir0zg/RD9dFOOckSRLFxcW1REdH/yctLe1znU63ZcSIEXUDBw6k2NhY2rNnD+3Zsyf68OHDuTt27Li0urr6yoqKipMPWoruNxHR8OHDv/jmm29G4mN7oQ3TxUNQZmbmvk2bNqlerzcoF/oZY6SqKpWVlUWUlZVdWVxcfKUsy7R582aKiIggnU5HbrebGhsbye12n/LtcH9oiaYoCjkcjo2iX40KFw6hE4Li4uK2G43G2oaGBuFvozpTYHi9XiIicrvP/3RBMAKHiCgyMpLi4uLWBat9gG6Nc84yMzO/oCBdG+mOpXfv3uWYidw9YJ5OCJJlmaempu6SZbyLqiMYY5SSkrI5KSmpNth9gfND6IQgVVVp0KBBm/CC8Y4xGo08Ojr6ryfu9gFAZxQUFLhSU1OrKQROXUK9pKSkHN20aRNOrQC6QpZlys3NXUsh8KMO5aIoCh8+fPhDXZlFDWLh9CpE+Xw+GjRo0Nro6IA9gxmWLBaLZ+rUqS/jrhVAAHz66afOpKSkEiKxjxR0p5KWlvai6JfZQ9dga4WwvLy8yqysrHcUBdOpzsTpdPIFCxa8IvJNiwBhb8mSJf1sNlsDhcCoIlSKJEmcMcZHjBjxDucct/gAAkmWZRozZsx7OL06tbhcrpbHHntsJC4gA2jgkUcemexyuTgRru0QEdfr9XzChAkP4VoOgEY457bJkyevlyQp6D/4UChDhw4tLi8vj+vyioWgwKGiG2CM1QwbNuzmuLi4i/57TnFxcZ5rrrnmVpfLdSzYfQEIa5Ik0dy5c5/U6/VBH2kEq1gsFn7TTTetCMT6BIAOqK2ttQ8aNKiAQiAARBedTsfz8vJWc84xfwBApJtuummoy+VqpBAIAhGFMcYVReHTp0/fzjm3Bmg1AkBHMcZozpw5C202W9ADQUTR6/V8ypQp31RVVSUGbCUCwIXhnMtz5sxZbTKZTo4GKAQCIlDFvzw6nY6PHz9+fVNTEwIHINg458YJEyY8L/r7UloXf+BYrVY+bty4Z3ANByCEcM4Nubm5D0dFRQU9LAJZEhIS3LNnz76Jc45XJwKEGs45mzlz5h+sVqubTowUutvpln/So16v5xkZGVvuuOOOXDzeABDi8vPzp6emph7sjqdbjDHucrma5s2b9xfOuUmbNQQAAbdx48b4yy+//EGn01mpKMopP2oKUqCcXtqPxE6EDR81atS7K1euvIRzrtd0BQGANpYuXdpr8uTJD2ZkZJQajcaQCh0i4hERETwxMbFp4sSJbz7xxBNTOOcSEdGyZcvwWA5Ad8UYo3379sVMnjz5qqysrH8mJyfXR0ZG8mA9NHoiaFqGDBnyxbx58+557bXXsvGJnYsTrtZdBCRJonXr1vV59dVXBxUWFo48duzYqIaGhoyqqqoIVVV1Ho8nIJ8D9tdhMBhIp9OR2Wx2WyyWA9HR0Vt79er1xRVXXPH1tddeuwefirm4IXQuQpxz9sorr/R+7rnnYrxe73CDwTCyvr4+rra2NtZms/Wtq6ujhoYG8vl856xHlmWKiooii8VCra2tJYyx4ri4uMampqb1ubm5ZbGxsTsWL15cJMuyF68UBb//AyJRJuw0L/mxAAAAAElFTkSuQmCC"
# ─────────────────────────────────────────────────────────────
# DGN BRAND CONSTANTS
# ─────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x11, 0x11, 0x11)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1E, 0x1E, 0x1E)
MID        = RGBColor(0x55, 0x55, 0x55)
LIGHT      = RGBColor(0x88, 0x88, 0x88)
SURFACE    = RGBColor(0xF2, 0xF2, 0xF2)
ULTRALIGHT = RGBColor(0xCC, 0xCC, 0xCC)

FONT = "Calibri"

def px(inches): return Inches(inches)
def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, px(x), px(y), px(w), px(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    return shape

def add_text(slide, text, x, y, w, h, size=13, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or MID
    return txBox

def add_label(slide, text, x, y, w):
    add_text(slide, text.upper(), x, y, w, 0.25, size=8, bold=True, color=LIGHT)

def add_divider(slide, x, y, w):
    line = slide.shapes.add_shape(1, px(x), px(y), px(w), px(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = ULTRALIGHT
    line.line.color.rgb = ULTRALIGHT

# ─────────────────────────────────────────────────────────────
# AI: GENERER TILBUDSTEKST
# ─────────────────────────────────────────────────────────────
def generer_tilbudstekst(input_data):
    client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))

    prompt = f"""Du er tilbudsassistent for Den Gode Nabo (DGN), et snekker/håndverksfirma i Drøbak.

Generer tilbudstekst basert på dette:
- Kunde: {input_data.get('kunde_navn', 'Kunde')}
- Prosjekttype: {input_data.get('prosjekt_type', 'Ukjent')}
- Beskrivelse: {input_data.get('beskrivelse', '')}
- Pris: {input_data.get('pris', '')} kr eks. mva.
- Tidsramme: {input_data.get('tidsramme', '')}
- Ekstra notater: {input_data.get('notater', '')}

Tone of Voice: Profesjonell men varm. Konkret og tillitsvekkende. Norsk håndverkstradisjon.

Svar KUN med JSON i dette formatet (ingen annen tekst):
{{
  "ingress": "2-3 setninger som åpner tilbudet varmt og konkret",
  "scope": "Hva jobben inkluderer, 2-4 punkter som streng med linjeskift",
  "materialer": "Materialvalg og kvalitet, 1-2 setninger",
  "garanti": "Kort setning om garanti/ettervern",
  "avslutning": "Vennlig avslutningssetning"
}}"""

    message = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=1000,
        messages=[{"role": "user", "content": prompt}]
    )

    raw = message.content[0].text.strip()
    raw = raw.replace("```json", "").replace("```", "").strip()
    return json.loads(raw)


# ─────────────────────────────────────────────────────────────
# PPTX BUILDER
# ─────────────────────────────────────────────────────────────
def bygg_tilbud_pptx(input_data, tekst):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]  # Blank layout

    kunde       = input_data.get("kunde_navn", "Kunde")
    prosjekt    = input_data.get("prosjekt_type", "Prosjekt")
    pris        = input_data.get("pris", "")
    tidsramme   = input_data.get("tidsramme", "")
    dato        = input_data.get("dato", "2026")
    adresse     = input_data.get("adresse", "")

    # ── SLIDE 1: COVER ──────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    add_rect(s1, 0, 0, 10, 5.625, WHITE)
    add_rect(s1, 4.8, 0, 5.2, 5.625, DARK)

    import io as _io
    _logo_bytes = base64.b64decode(DGN_LOGO_B64)
    _logo_stream = _io.BytesIO(_logo_bytes)
    s1.shapes.add_picture(_logo_stream, px(0.6), px(1.2), px(2.2), px(2.24))
    add_text(s1, prosjekt, 0.65, 3.7, 3.8, 0.4,
             size=13, color=LIGHT)
    add_text(s1, dato, 0.65, 5.1, 3.5, 0.3,
             size=8, bold=True, color=LIGHT)

    # ── SLIDE 2: PROSJEKTDATA ────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    add_rect(s2, 0, 0, 10, 5.625, WHITE)

    add_text(s2, "Prosjektdata", 0.6, 2.0, 3.8, 1.0,
             size=32, color=BLACK)

    fields = [
        ("KUNDE",      kunde),
        ("PROSJEKT",   prosjekt),
        ("ADRESSE",    adresse or "—"),
        ("DATO",       dato),
    ]
    yy = 0.5
    for label, value in fields:
        add_divider(s2, 4.8, yy, 5.0)
        yy += 0.15
        add_label(s2, label, 4.8, yy, 5.0)
        yy += 0.27
        add_text(s2, value, 4.8, yy, 5.0, 0.3, size=12, color=MID)
        yy += 0.45
    add_divider(s2, 4.8, yy, 5.0)

    # ── SLIDE 3: INGRESS ─────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    add_rect(s3, 0, 0, 10, 5.625, WHITE)

    add_text(s3, "Tilbud", 0.6, 0.7, 9.0, 0.7,
             size=32, color=BLACK)
    add_divider(s3, 0.6, 1.55, 8.8)
    add_text(s3, tekst.get("ingress", ""), 0.6, 1.75, 8.8, 2.0,
             size=13, color=MID)

    # ── SLIDE 4: OMFANG ──────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    add_rect(s4, 0, 0, 10, 5.625, WHITE)
    add_rect(s4, 0, 0, 4.4, 5.625, SURFACE)

    add_label(s4, "OMFANG", 4.8, 0.45, 5.0)
    add_text(s4, "Hva jobben inkluderer", 4.8, 0.75, 5.0, 0.6,
             size=22, color=BLACK)

    scope_lines = tekst.get("scope", "").split("\n")
    yy = 1.55
    for line in scope_lines:
        if line.strip():
            add_text(s4, f"— {line.strip()}", 4.8, yy, 4.9, 0.35,
                     size=12, color=MID)
            yy += 0.4

    add_text(s4, tekst.get("materialer", ""), 0.4, 2.2, 3.6, 1.5,
             size=11, color=LIGHT, italic=True)

    # ── SLIDE 5: PRIS ────────────────────────────────────────
    s5 = prs.slides.add_slide(blank)
    add_rect(s5, 0, 0, 10, 5.625, WHITE)
    add_rect(s5, 5.8, 0, 4.2, 5.625, SURFACE)

    add_text(s5, "Pris og\nbetingelser", 0.6, 1.8, 4.8, 1.5,
             size=32, color=BLACK)

    price_items = [
        ("Fastpris",      f"{pris} kr eks. mva." if pris else "—"),
        ("Tidsramme",     tidsramme or "—"),
        ("Faktura",       "100% ved leveranse"),
        ("Garanti",       tekst.get("garanti", "—")),
    ]
    yy = 0.6
    for label, value in price_items:
        add_text(s5, label, 6.1, yy, 3.7, 0.28,
                 size=11, bold=True, color=BLACK)
        yy += 0.3
        add_text(s5, value, 6.1, yy, 3.7, 0.3,
                 size=11, color=MID)
        yy += 0.65

    # ── SLIDE 6: AVSLUTNING ──────────────────────────────────
    s6 = prs.slides.add_slide(blank)
    add_rect(s6, 0, 0, 10, 5.625, DARK)

    add_text(s6, tekst.get("avslutning", "Takk for oppdraget."),
             1.0, 2.0, 8.0, 1.0,
             size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s6, "dengodanabo.no",
             1.0, 4.9, 8.0, 0.4,
             size=9, color=LIGHT, align=PP_ALIGN.CENTER)

    # ── SKRIV TIL BUFFER ─────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────────────────────
# ENDEPUNKT: /generer-tilbud
# ─────────────────────────────────────────────────────────────
@app.route('/generer-tilbud', methods=['POST'])
def generer_tilbud():
    """
    Input (JSON):
      kunde_navn    str   — "Fru Hansen"
      prosjekt_type str   — "Platting", "Tilbygg", "Carport" etc.
      beskrivelse   str   — fritekst fra befaring
      pris          str   — "45000"
      tidsramme     str   — "3-4 uker"
      dato          str   — "Mars 2026"
      adresse       str   — "Storgata 12, Drøbak"
      notater       str   — ekstra stikkord

    Output:
      .pptx fil som nedlasting
    """
    data = request.json
    if not data:
        return jsonify({"error": "Mangler JSON-body"}), 400

    required = ["kunde_navn", "prosjekt_type", "beskrivelse"]
    for field in required:
        if not data.get(field):
            return jsonify({"error": f"Mangler felt: {field}"}), 400

    try:
        # 1. Generer tekst med Claude
        tekst = generer_tilbudstekst(data)

        # 2. Bygg PPTX
        pptx_buf = bygg_tilbud_pptx(data, tekst)

        # 3. Returner filen
        filnavn = f"DGN_Tilbud_{data['kunde_navn'].replace(' ', '_')}.pptx"
        return send_file(
            pptx_buf,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=filnavn
        )

    except json.JSONDecodeError as e:
        return jsonify({"error": f"AI returnerte ugyldig JSON: {str(e)}"}), 500
    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


# ─────────────────────────────────────────────────────────────
# EKSISTERENDE ENDEPUNKTER (urørt)
# ─────────────────────────────────────────────────────────────
@app.route('/visualiser_prosjekt', methods=['POST'])
def visualiser():
    data = request.json
    beskrivelse = data.get("beskrivelse")
    foto_base64 = data.get("foto_base64")
    foto_url = data.get("foto_url")

    if not foto_base64 and not foto_url:
        return jsonify({"error": "Mangler bilde — send enten foto_base64 eller foto_url"}), 400

    try:
        if foto_base64:
            img_bytes = base64.b64decode(foto_base64)
        else:
            response = requests.get(foto_url, timeout=10)
            response.raise_for_status()
            img_bytes = response.content

        img = Image.open(io.BytesIO(img_bytes))
        img.thumbnail((1024, 1024))
        buffered = io.BytesIO()
        img.save(buffered, format="JPEG")
        img_b64 = base64.b64encode(buffered.getvalue()).decode()

        prompt = f"""Generer en fotorealistisk visualisering av dette byggeprosjektet etter ferdigstillelse:

{beskrivelse}

Behold samme perspektiv og kameravinkel som i befaringsbildet.
Resultatet skal se ut som et ekte fotografi."""

        model = genai.GenerativeModel('gemini-2.0-flash-exp-image-generation')
        res = model.generate_content(
            [{"mime_type": "image/jpeg", "data": img_b64}, prompt],
            generation_config={"response_modalities": ["TEXT", "IMAGE"]}
        )

        for part in res.candidates[0].content.parts:
            if hasattr(part, 'inline_data') and part.inline_data is not None:
                raw = part.inline_data.data
                mime = part.inline_data.mime_type
                if isinstance(raw, bytes):
                    img_out = base64.b64encode(raw).decode("utf-8")
                elif isinstance(raw, str):
                    img_out = raw
                else:
                    continue
                if img_out:
                    return jsonify({"visualisering_base64": img_out, "mime_type": mime})

        parts_debug = []
        for p in res.candidates[0].content.parts:
            if hasattr(p, 'inline_data') and p.inline_data:
                parts_debug.append({"type": str(type(p.inline_data.data)), "len": len(p.inline_data.data) if p.inline_data.data else 0, "mime": p.inline_data.mime_type})
            elif hasattr(p, 'text'):
                parts_debug.append({"text": p.text[:200]})
        return jsonify({"error": "Ingen bilde i respons", "debug": parts_debug}), 500

    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "trace": traceback.format_exc()}), 500


@app.route('/sse')
def sse():
    return jsonify({
        "tools": [{
            "name": "visualiser_prosjekt",
            "description": "Tar inn befaringsfoto og prosjektbeskrivelse, returnerer fotorealistisk visualisering av ferdig prosjekt.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "foto_base64": {"type": "string", "description": "Base64-kodet befaringsfoto uten data:image/jpeg;base64,-prefix"},
                    "beskrivelse": {"type": "string", "description": "Hva som skal gjøres, f.eks. 'ny terrasse i trykkimpregnert tre'"}
                },
                "required": ["foto_base64", "beskrivelse"]
            }
        }]
    })


@app.route('/health')
def health():
    return jsonify({"status": "ok"})


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 8080)))
